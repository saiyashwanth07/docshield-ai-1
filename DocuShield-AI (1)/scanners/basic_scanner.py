import os
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
import openpyxl
from odf.opendocument import load
from odf.text import P
from io import StringIO

# Simple keyword-based scanner
MALWARE_KEYWORDS = ["malware", "virus", "trojan", "ransomware", "exploit", "phishing", "keylogger", "worm"]

def extract_text_from_file(file_path):
    ext = os.path.splitext(file_path)[-1].lower()

    try:
        if ext == '.pdf':
            reader = PdfReader(file_path)
            return "\n".join(page.extract_text() for page in reader.pages if page.extract_text())
        
        elif ext == '.docx':
            doc = Document(file_path)
            return "\n".join([para.text for para in doc.paragraphs])
        
        elif ext == '.pptx':
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        
        elif ext == '.xlsx':
            wb = openpyxl.load_workbook(file_path)
            text = ""
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    text += " ".join([str(cell) for cell in row if cell]) + "\n"
            return text
        
        elif ext == '.odt':
            textdoc = load(file_path)
            allparas = textdoc.getElementsByType(P)
            return "\n".join([str(p) for p in allparas])
        
        elif ext == '.txt':
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()

        elif ext == '.rtf':
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()

        else:
            return ""
    
    except Exception as e:
        return f"Error extracting text: {e}"

def scan_document(file_path):
    content = extract_text_from_file(file_path).lower()
    if not content.strip():
        return "Unable to extract readable content from this file."

    found_keywords = [kw for kw in MALWARE_KEYWORDS if kw in content]
    
    if found_keywords:
        return f"⚠️ Potential threats detected: {', '.join(found_keywords)}"
    else:
        return "✅ No known malware indicators found in the document."
